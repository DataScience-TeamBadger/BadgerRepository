from pptx import Presentation

#Badger Data Science Presentation Code
#Presentation Conducted -
#Data Science - Alan Jamieson - COSC 480, Spring 2017

#Title Slide
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs. slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Maximizing Spending Allocation on Metropolitan Mass Transportation"
subtitle.text = "Badger Data Science"

#Slide 1 - Basic Information
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'What is Metropolitan Mass Transportation?'

tf = body_shape.text_frame
tf.text = 'Definition: The transportation of large numbers of people by means of buses, subways trains, etc., especially within urban areas' \
          '\n - Merriam-Webster'

tf = body_shape.text_frame
tf.text = 'What does it consist of?'

p = tf.add_paragraph()
p.text = 'How many people are being moved.'
p.level = 1

p = tf.add_paragraph()
p.text = 'How far do these people have to be moved.'
p.level = 1

p = tf.add_paragraph()
p.text = 'How much does it cost to move these people.'
p.level = 1

tf = body_shape.text_frame
tf.text = 'Goal: To Maximize Transportation Budget Spending by Predicting the best Allocation of Budget to Metro and Bus, in order to move the most people, the furthest distance'

#Slide 2 - Why Is it Important
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'What makes Mass Transit an Important Issue?'

tf = body_shape.text_frame
tf.text = 'Moving to Equity - Harvard University, The Civil Rights Project'
#Most Americans rely on their car for transportation to work.
p = tf.add_paragraph()
p.text = 'Minorities have a much higher rate of not owning a form of personal transportation.'
p.level = 1

p = tf.add_paragraph()
p.text = 'Ownership of personal transportation modes per household:\n White: 97% \n African American: 76%\n  Latino: 83%\n Asian American: 13%'
 ##Want to add graph from justification paper

#Slide 3 - Why Is it Important pt2
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Theory Behind the Issue -'

tf = body_shape.text_frame
p = tf.add_paragraph()
p.text = 'Spacial Mismatch Hypothesis - John Kain(1968)'
p.level = 1

p = tf.add_paragraph()
p.text = '"those who most need entry-level jobs (primarily people of color) generally live in central cities while entry-level jobs are mostly in suburban locations that are not easily accessible from central cities." '
p.level = 2

p = tf.add_paragraph()
p.text = 'Social Exclusion - "What can happen when people or areas suffer from a combination of linked problems such as unemployment, poor skills, low incomes, poor housing, high crime, bad health and family breakdown" \n - Enlgish Government'
p.level = 1
prs.save('BadgerDataScience_Presentation_03.pptx')
